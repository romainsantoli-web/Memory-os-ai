# Fichier : memory_os_ai.py
# Description : Script principal pour Memory OS AI, un système d'IA avancé pour la gestion et l'analyse de documents localement.
# Copyright (c) 2025 Kocupyr Romain
# Licence : LGPL v3 pour usage non commercial ; licence commerciale payante pour usage commercial.

import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from llama_cpp import Llama
import time
import fitz  # PyMuPDF pour PDFs
import os
import pickle
from multiprocessing import Pool
import threading
import torch
from functools import lru_cache
import re
import subprocess
from docx import Document  # Pour .docx
import textract  # Pour .doc et .ppt
from PIL import Image  # Pour traitement d’images
import pytesseract  # Pour OCR
from pptx import Presentation  # Pour .pptx
import shutil  # Pour copier les fichiers
import whisper  # Pour transcription audio avec Whisper v2
from huggingface_hub import snapshot_download, hf_hub_download  # Pour télécharger les modèles
import requests  # Pour télécharger Mistral-7B

# Désactiver parallélisme tokenizers
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ["PYTORCH_CUDA_ALLOC_CONF"] = "expandable_segments:True"
torch.backends.cuda.matmul.allow_tf32 = True

# --- Configuration globale ---
# Détection automatique CPU/GPU
DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
print(f"[INFO] Utilisation du dispositif : {DEVICE}")

# Chemins relatifs au dossier du script
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
MODELE_SEMANTIQUE_PATH = os.path.join(BASE_DIR, "emb", "all-MiniLM-L6-v2")
CHEMIN_MODELE = os.path.join(BASE_DIR, "mistral-7b", "mistral-7b-instruct-v0.2.Q4_K_M.gguf")
WHISPER_MODEL_PATH = os.path.join(BASE_DIR, "emb", "whisper_base.pt")
CACHE_EMBEDDINGS = os.path.join(BASE_DIR, "embeddings_cache.pkl")
DOSSIER_FICHIERS = os.path.join(BASE_DIR, "pdfs")
EXTENSIONS_ACCEPTEES = {".pdf", ".txt", ".docx", ".doc", ".png", ".jpeg", ".ppt", ".pptx", ".mp3", ".wav", ".ogg", ".flac"}

# Créer les dossiers si nécessaire
os.makedirs(os.path.join(BASE_DIR, "emb"), exist_ok=True)
os.makedirs(os.path.join(BASE_DIR, "mistral-7b"), exist_ok=True)
os.makedirs(os.path.join(BASE_DIR, "pdfs"), exist_ok=True)

# Télécharger les modèles si nécessaire
def download_models():
    """Télécharge les modèles nécessaires si les fichiers ou dossiers n'existent pas."""
    # SentenceTransformers (all-MiniLM-L6-v2)
    if not os.path.exists(MODELE_SEMANTIQUE_PATH):
        print("[TÉLÉCHARGEMENT] Téléchargement du modèle SentenceTransformers (all-MiniLM-L6-v2)...")
        snapshot_download(repo_id="sentence-transformers/all-MiniLM-L6-v2", local_dir=MODELE_SEMANTIQUE_PATH)

    # Mistral-7B (mistral-7b-instruct-v0.2.Q4_K_M.gguf)
    if not os.path.exists(CHEMIN_MODELE):
        print("[TÉLÉCHARGEMENT] Téléchargement du modèle Mistral-7B (mistral-7b-instruct-v0.2.Q4_K_M.gguf)...")
        url = "https://huggingface.co/TheBloke/Mixtral-8x7B-Instruct-v0.1-GGUF/resolve/main/mixtral-8x7b-instruct-v0.1.Q4_K_M.gguf"
        response = requests.get(url, stream=True)
        with open(CHEMIN_MODELE, "wb") as f:
            for chunk in response.iter_content(chunk_size=8192):
                if chunk:
                    f.write(chunk)

    # Whisper (base)
    if not os.path.exists(WHISPER_MODEL_PATH):
        print("[TÉLÉCHARGEMENT] Téléchargement du modèle Whisper (base)...")
        hf_hub_download(repo_id="openai/whisper-base", filename="model.pt", local_dir=os.path.join(BASE_DIR, "emb"))
        os.rename(os.path.join(BASE_DIR, "emb", "model.pt"), WHISPER_MODEL_PATH)

# Télécharger les modèles au démarrage
download_models()

# Initialisation des modèles
print("[DÉBUT] Chargement de SentenceTransformer...")
MODELE_SEMANTIQUE = SentenceTransformer(MODELE_SEMANTIQUE_PATH, device=DEVICE)

print("[DÉBUT] Chargement de Mistral (une seule fois)...")
MODELE_GENERATIF = Llama(model_path=CHEMIN_MODELE, n_ctx=2048, n_gpu_layers=32 if DEVICE == "cuda" else 0, n_threads=8)

print("[DÉBUT] Chargement de Whisper v2 (base)...")
if not os.path.exists(WHISPER_MODEL_PATH):
    print(f"Erreur : Le modèle Whisper v2 (base) n’est pas trouvé à {WHISPER_MODEL_PATH}.")
    exit(1)
original_torch_load = torch.load
def safe_torch_load(*args, **kwargs):
    kwargs["weights_only"] = True
    return original_torch_load(*args, **kwargs)
torch.load = safe_torch_load
MODELE_WHISPER = whisper.load_model(WHISPER_MODEL_PATH)
torch.load = original_torch_load

# FAISS (CPU ou GPU selon la disponibilité)
RES_FAISS = faiss.StandardGpuResources() if DEVICE == "cuda" else None
DIMENSION = 384
TEXTES_DOCUMENTS = []
INDEX_FAISS = None
LIENS_PRECALCULES = None
FICHIERS_SEGMENTS_MAP = {}
FICHIERS_RESUMES = {}
LOCK = threading.Lock()
SEMAPHORE = threading.Semaphore(3)
MODEL_LOCK = threading.Lock()

# --- Fonctions de traitement ---
def process_file(fichier, dossier):
    """Traite un fichier et extrait son contenu sous forme de segments."""
    chemin_fichier = os.path.join(dossier, fichier)
    try:
        if fichier.endswith(".pdf"):
            texte = ""
            nb_pages = 0
            with fitz.open(chemin_fichier) as doc:
                nb_pages = len(doc)
                for page in doc:
                    texte_page = page.get_text()
                    if texte_page.strip():
                        texte += texte_page
                    else:
                        pix = page.get_pixmap()
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        texte += pytesseract.image_to_string(img, lang="eng+fra", config="--psm 3")
            texte = texte.strip()

        elif fichier.endswith(".txt"):
            with open(chemin_fichier, "r", encoding="utf-8") as f:
                texte = f.read()
                nb_pages = 1

        elif fichier.endswith(".docx"):
            doc = Document(chemin_fichier)
            texte = "\n".join(para.text for para in doc.paragraphs)
            nb_pages = 1

        elif fichier.endswith(".doc"):
            texte = textract.process(chemin_fichier, method="antiword").decode("utf-8")
            nb_pages = 1

        elif fichier.endswith((".png", ".jpeg")):
            img = Image.open(chemin_fichier).convert("L")
            texte = pytesseract.image_to_string(img, lang="eng+fra", config="--psm 3")
            nb_pages = 1

        elif fichier.endswith(".pptx"):
            prs = Presentation(chemin_fichier)
            texte = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texte += shape.text + "\n"
            texte = texte.strip()
            nb_pages = len(prs.slides)

        elif fichier.endswith(".ppt"):
            texte = textract.process(chemin_fichier).decode("utf-8")
            nb_pages = 1

        elif fichier.endswith((".mp3", ".wav", ".ogg", ".flac")):
            result = MODELE_WHISPER.transcribe(chemin_fichier, language="fr")
            texte = result["text"].strip()
            nb_pages = 1

        else:
            raise ValueError(f"Format non pris en charge pour {fichier}")
    except Exception as e:
        print(f"Erreur lors du traitement de {fichier} : {e}")
        return [], fichier, 0
    
    segments = [texte[i:i+128] for i in range(0, len(texte), 100)]
    return segments, fichier, nb_pages

def charger_fichiers(dossier):
    """Charge les fichiers depuis un dossier et les segmente."""
    global FICHIERS_SEGMENTS_MAP
    fichiers = [f for f in os.listdir(dossier) if os.path.splitext(f)[1].lower() in EXTENSIONS_ACCEPTEES]
    with Pool(processes=12) as pool:
        results = pool.starmap(process_file, [(f, dossier) for f in fichiers])
    segments, fichiers_segments, total_pages = [], {}, 0
    debut_idx = 0
    for segments_locaux, fichier, nb_pages in results:
        if segments_locaux:
            segments.extend(segments_locaux)
            fichiers_segments[fichier] = segments_locaux
            FICHIERS_SEGMENTS_MAP[fichier] = (debut_idx, debut_idx + len(segments_locaux))
            debut_idx += len(segments_locaux)
            total_pages += nb_pages
            print(f"[FICHIER] {fichier} chargé : {len(segments_locaux)} segments, {nb_pages} pages")
    return segments, fichiers_segments, total_pages

@lru_cache(maxsize=1024)
def encode_chunk_cached(chunk_tuple):
    """Encode un chunk de texte en embeddings (mis en cache)."""
    chunk = list(chunk_tuple)
    with SEMAPHORE:
        embeddings = MODELE_SEMANTIQUE.encode(chunk, batch_size=4096, convert_to_numpy=True, normalize_embeddings=True)
    return embeddings

def creer_embeddings_multidimensionnels(textes):
    """Crée des embeddings multidimensionnels pour les textes."""
    debut = time.time()
    if DEVICE == "cuda":
        torch.cuda.empty_cache()
    embeddings = MODELE_SEMANTIQUE.encode(textes, batch_size=4096, convert_to_numpy=True, normalize_embeddings=True)
    if DEVICE == "cuda":
        torch.cuda.synchronize()
    temps = time.time() - debut
    print(f"[EMBED] {len(textes)} textes encodés en {temps:.2f}s")
    return embeddings.astype("float32")

def ingerer_textes(textes, fichiers_segments, total_pages):
    """Ingère les textes et crée un index FAISS."""
    global INDEX_FAISS, TEXTES_DOCUMENTS
    debut = time.time()
    
    if os.path.exists(CACHE_EMBEDDINGS):
        with open(CACHE_EMBEDDINGS, "rb") as f:
            embeddings = pickle.load(f)
        print(f"[CACHE] Embeddings chargés depuis {CACHE_EMBEDDINGS}")
    else:
        embeddings = creer_embeddings_multidimensionnels(textes)
        with open(CACHE_EMBEDDINGS, "wb") as f:
            pickle.dump(embeddings, f)
        print(f"[CACHE] Embeddings sauvegardés dans {CACHE_EMBEDDINGS}")
    
    if DEVICE == "cuda":
        INDEX_FAISS = faiss.index_cpu_to_gpu(RES_FAISS, 0, faiss.IndexFlatL2(DIMENSION))
    else:
        INDEX_FAISS = faiss.IndexFlatL2(DIMENSION)
    INDEX_FAISS.add(embeddings)
    TEXTES_DOCUMENTS.extend(textes)
    
    temps_total = time.time() - debut
    mots_total = sum(len(t.split()) for t in textes)
    nb_fichiers = len(fichiers_segments)
    print(f"[INGESTION] {nb_fichiers} fichiers, {total_pages} pages, {mots_total} mots, temps : {temps_total:.2f}s")
    if DEVICE == "cuda":
        torch.cuda.empty_cache()
    return embeddings, fichiers_segments

def resumer_fichiers(fichiers_segments):
    """Résume les fichiers chargés."""
    global FICHIERS_RESUMES
    resumes, debut_total = {}, time.time()
    for fichier, segments in fichiers_segments.items():
        debut = time.time()
        if DEVICE == "cuda":
            torch.cuda.empty_cache()
        contexte = "\n".join(segments)[:2048]
        prompt = f"[INST] Résume le contenu du fichier suivant en quelques phrases en français :\nContexte : {contexte} [/INST]"
        with MODEL_LOCK:
            resume = MODELE_GENERATIF(prompt, max_tokens=400, temperature=0.3)["choices"][0]["text"]
        resumes[fichier] = resume
        FICHIERS_RESUMES[fichier] = resume
        print(f"[RÉSUMÉ] {fichier} : {resume} (temps : {time.time() - debut:.2f}s)")
    
    print(f"[RÉSUMÉS] Temps total : {time.time() - debut_total:.2f}s")
    return resumes

def rechercher_occurrences_faiss(mot_cle, fichiers_segments):
    """Compte les occurrences d’un mot-clé dans les documents via l’index FAISS."""
    emb_mot_cle = encode_chunk_cached(tuple([mot_cle]))[0]
    distances, indices = INDEX_FAISS.search(np.array([emb_mot_cle]), k=500)
    
    resultats = {}
    total_occurrences = 0
    
    for fichier, (debut_idx, fin_idx) in FICHIERS_SEGMENTS_MAP.items():
        segments_pertinents = [TEXTES_DOCUMENTS[idx] for i, idx in enumerate(indices[0]) if debut_idx <= idx < fin_idx and distances[0][i] < 0.8]
        if segments_pertinents:
            texte_complet = " ".join(segments_pertinents)
            occurrences = len(re.findall(r'\b' + re.escape(mot_cle) + r'\b', texte_complet, re.IGNORECASE))
            if occurrences > 0:
                resultats[fichier] = occurrences
                total_occurrences += occurrences
    
    resultats_trie = dict(sorted(resultats.items(), key=lambda x: x[1], reverse=True))
    return total_occurrences, resultats_trie

def generate_dynamic_response(prompt, initial_max_tokens=600):
    """Génère une réponse dynamique avec le modèle génératif."""
    with MODEL_LOCK:
        response = MODELE_GENERATIF(prompt, max_tokens=initial_max_tokens, temperature=0.3)["choices"][0]["text"]
    return response

def main():
    """Fonction principale pour exécuter le script."""
    print("Memory OS AI - Copyright (c) 2025 [Ton Nom]")
    print("Ce logiciel est sous licence LGPL v3 pour une utilisation non commerciale.")
    print("Pour une utilisation commerciale, veuillez acheter une licence : [ton email]")
    print("Soutenez le projet en faisant un don : [Lien vers Patreon/Open Collective]")
    print("---")

    debut_global = time.time()
    print(f"[DÉBUT] Chargement des fichiers depuis {DOSSIER_FICHIERS}")
    textes_fichiers, fichiers_segments, total_pages = charger_fichiers(DOSSIER_FICHIERS)
    embeddings, fichiers_segments = ingerer_textes(textes_fichiers, fichiers_segments, total_pages)
    resumer_fichiers(fichiers_segments)

    # Interface en ligne de commande simple
    print("\nBienvenue dans Memory OS AI ! Tapez 'exit' pour quitter.")
    while True:
        requete = input("Entrez votre requête (ex. 'recherche qubit', 'rapport sur X') : ").strip()
        if requete.lower() == "exit":
            break
        
        debut = time.time()
        if "recherche" in requete.lower():
            mot_cle = requete.split("recherche", 1)[1].strip()
            total_occurrences, resultats = rechercher_occurrences_faiss(mot_cle, fichiers_segments)
            if resultats:
                resumes_pertinents = {fichier: FICHIERS_RESUMES[fichier] for fichier in resultats.keys()}
                prompt = f"[INST] Réponds en français uniquement. Fournis une réponse claire et bien présentée en français, basée sur les résultats suivants, sans mentionner 'segments'. Indique le nombre total d’occurrences du mot-clé '{mot_cle}', liste les fichiers avec leur nombre d’occurrences, triés par pertinence (nombre d’occurrences décroissant), et résume brièvement chaque fichier où le mot-clé apparaît :\nRésultats : Total = {total_occurrences}, Fichiers = {resultats}, Résumés = {resumes_pertinents}\nRéponse (en français uniquement) : [/INST]"
                reponse = generate_dynamic_response(prompt)
                print(f"Réponse : {reponse} (temps : {time.time() - debut:.2f}s)")
            else:
                print(f"Aucune occurrence du mot-clé '{mot_cle}' n’a été trouvée dans les fichiers.")
        
        elif "rapport" in requete.lower():
            sujet = requete.split("rapport", 1)[1].strip()
            emb_requete = encode_chunk_cached(tuple([sujet]))[0]
            distances, indices = INDEX_FAISS.search(np.array([emb_requete]), k=200)
            contexte_general = "\n".join([TEXTES_DOCUMENTS[idx] for idx in indices[0] if idx < len(TEXTES_DOCUMENTS)])[:1500]
            prompt = f"[INST] Réponds en français uniquement. Crée un rapport structuré (introduction, contenu principal, conclusion) basé uniquement sur le contenu des fichiers extraits de l’index FAISS concernant le sujet '{sujet}'. Ne génère rien hors sujet, ne mentionne pas 'segments', et n’ajoute pas d’excuses :\nContenu des fichiers : {contexte_general}\nListe des fichiers : {', '.join(fichiers_segments.keys())}\nRapport (en français uniquement) : [/INST]"
            reponse = generate_dynamic_response(prompt, initial_max_tokens=1000)
            print(f"Réponse : {reponse} (temps : {time.time() - debut:.2f}s)")
        
        else:
            print("Commande non reconnue. Essayez 'recherche <mot-clé>' ou 'rapport <sujet>'.")

    print(f"[GLOBAL] Temps total : {time.time() - debut_global:.2f}s")

if __name__ == "__main__":
    main()
