"""Core document processing engine for Memory OS AI.

Handles document ingestion, embedding, FAISS indexing, and search.
No LLM dependency — all generation is delegated to the caller (any AI model via MCP).
"""

from __future__ import annotations

import logging
import os
import re
import threading
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

import numpy as np

logger = logging.getLogger("memory_os_ai")

# --- Lazy imports for heavy libraries ---
# These are imported at first use to speed up module load time
_faiss = None
_SentenceTransformer = None
_torch = None


def _ensure_faiss():
    global _faiss
    if _faiss is None:
        import faiss
        _faiss = faiss
    return _faiss


def _ensure_torch():
    global _torch
    if _torch is None:
        import torch
        _torch = torch
    return _torch


def _ensure_sentence_transformer():
    global _SentenceTransformer
    if _SentenceTransformer is None:
        from sentence_transformers import SentenceTransformer
        _SentenceTransformer = SentenceTransformer
    return _SentenceTransformer


# --- Document processing imports (lazy) ---
def _extract_pdf(path: str) -> tuple[str, int]:
    """Extract text from PDF, with OCR fallback."""
    import fitz
    from PIL import Image
    import pytesseract

    text = ""
    with fitz.open(path) as doc:
        nb_pages = len(doc)
        for page in doc:
            page_text = page.get_text()
            if page_text.strip():
                text += page_text
            else:
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text += pytesseract.image_to_string(img, lang="eng+fra", config="--psm 3")
    return text.strip(), nb_pages


def _extract_txt(path: str) -> tuple[str, int]:
    with open(path, "r", encoding="utf-8") as f:
        return f.read(), 1


def _extract_docx(path: str) -> tuple[str, int]:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs), 1


def _extract_doc(path: str) -> tuple[str, int]:
    import textract
    return textract.process(path, method="antiword").decode("utf-8"), 1


def _extract_image(path: str) -> tuple[str, int]:
    from PIL import Image
    import pytesseract
    img = Image.open(path).convert("L")
    return pytesseract.image_to_string(img, lang="eng+fra", config="--psm 3"), 1


def _extract_pptx(path: str) -> tuple[str, int]:
    from pptx import Presentation
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip(), len(prs.slides)


def _extract_ppt(path: str) -> tuple[str, int]:
    import textract
    return textract.process(path).decode("utf-8"), 1


def _extract_audio(path: str, language: str = "fr") -> tuple[str, int]:
    """Transcribe audio using Whisper."""
    import whisper
    torch = _ensure_torch()

    # Safe torch.load wrapper
    original_load = torch.load
    def safe_load(*args, **kwargs):
        kwargs["weights_only"] = True
        return original_load(*args, **kwargs)
    torch.load = safe_load

    try:
        model = whisper.load_model("base")
        result = model.transcribe(path, language=language)
        return result["text"].strip(), 1
    finally:
        torch.load = original_load


# Extractors registry
EXTRACTORS: dict[str, callable] = {
    ".pdf": _extract_pdf,
    ".txt": _extract_txt,
    ".docx": _extract_docx,
    ".doc": _extract_doc,
    ".png": _extract_image,
    ".jpeg": _extract_image,
    ".jpg": _extract_image,
    ".pptx": _extract_pptx,
    ".ppt": _extract_ppt,
    ".mp3": _extract_audio,
    ".wav": _extract_audio,
    ".ogg": _extract_audio,
    ".flac": _extract_audio,
}

SUPPORTED_EXTENSIONS = set(EXTRACTORS.keys())


# --- Data classes ---
@dataclass
class DocumentInfo:
    """Metadata about an indexed document."""
    filename: str
    nb_pages: int
    nb_segments: int
    nb_words: int
    segment_start: int
    segment_end: int


@dataclass
class SearchResult:
    """A single search result."""
    filename: str
    segment_text: str
    distance: float
    segment_index: int


@dataclass
class OccurrenceResult:
    """Keyword occurrence count per file."""
    total: int
    by_file: dict[str, int]


# --- Memory Engine ---
class MemoryEngine:
    """Core engine for document ingestion and semantic search.

    This class manages the FAISS index, embeddings, and document metadata.
    It does NOT include any LLM — generation is delegated to the caller via MCP.
    """

    def __init__(
        self,
        model_name: str = "all-MiniLM-L6-v2",
        dimension: int = 384,
        segment_size: int = 128,
        segment_overlap: int = 28,
        cache_dir: Optional[str] = None,
    ):
        self.model_name = model_name
        self.dimension = dimension
        self.segment_size = segment_size
        self.segment_overlap = segment_overlap
        self.cache_dir = cache_dir

        # State
        self._documents: dict[str, DocumentInfo] = {}
        self._segments: list[str] = []
        self._index = None  # FAISS index
        self._model = None  # SentenceTransformer
        self._device: Optional[str] = None
        self._lock = threading.Lock()
        self._initialized = False

    @property
    def device(self) -> str:
        if self._device is None:
            torch = _ensure_torch()
            self._device = "cuda" if torch.cuda.is_available() else "cpu"
        return self._device

    @property
    def model(self):
        if self._model is None:
            ST = _ensure_sentence_transformer()
            model_path = self.model_name
            if self.cache_dir:
                local_path = os.path.join(self.cache_dir, self.model_name)
                if os.path.isdir(local_path):
                    model_path = local_path
            self._model = ST(model_path, device=self.device)
        return self._model

    @property
    def is_initialized(self) -> bool:
        return self._initialized and self._index is not None

    @property
    def document_count(self) -> int:
        return len(self._documents)

    @property
    def segment_count(self) -> int:
        return len(self._segments)

    def _segment_text(self, text: str) -> list[str]:
        """Split text into overlapping segments."""
        step = self.segment_size - self.segment_overlap
        if step <= 0:
            step = self.segment_size
        return [text[i:i + self.segment_size] for i in range(0, len(text), step)]

    def _encode(self, texts: list[str]) -> np.ndarray:
        """Encode texts to embeddings."""
        embeddings = self.model.encode(
            texts,
            batch_size=256,
            convert_to_numpy=True,
            normalize_embeddings=True,
            show_progress_bar=False,
        )
        return embeddings.astype("float32")

    def _build_index(self, embeddings: np.ndarray) -> None:
        """Build or rebuild the FAISS index."""
        faiss = _ensure_faiss()

        if self.device == "cuda":
            res = faiss.StandardGpuResources()
            self._index = faiss.index_cpu_to_gpu(res, 0, faiss.IndexFlatL2(self.dimension))
        else:
            self._index = faiss.IndexFlatL2(self.dimension)

        self._index.add(embeddings)

    def ingest(
        self,
        folder_path: str,
        extensions: Optional[set[str]] = None,
        force_reindex: bool = False,
    ) -> dict:
        """Ingest documents from a folder.

        Returns a summary dict with counts and timing.
        """
        if extensions is None:
            extensions = SUPPORTED_EXTENSIONS
        else:
            extensions = {e.lower() for e in extensions}

        if not os.path.isdir(folder_path):
            return {"ok": False, "error": f"Folder not found: {folder_path}"}

        cache_path = os.path.join(folder_path, "embeddings_cache.npy") if not self.cache_dir else os.path.join(self.cache_dir, "embeddings_cache.npy")

        start = time.time()
        all_segments: list[str] = []
        documents: dict[str, DocumentInfo] = {}
        total_pages = 0

        # Process files
        files = sorted(
            f for f in os.listdir(folder_path)
            if os.path.splitext(f)[1].lower() in extensions
        )

        results = []
        for filename in files:
            ext = os.path.splitext(filename)[1].lower()
            filepath = os.path.join(folder_path, filename)
            extractor = EXTRACTORS.get(ext)
            if not extractor:
                continue
            try:
                text, nb_pages = extractor(filepath)
                segments = self._segment_text(text)
                results.append((filename, segments, nb_pages, len(text.split())))
            except Exception as e:
                results.append((filename, [], 0, 0))

        # Build segments map
        offset = 0
        for filename, segments, nb_pages, nb_words in results:
            if segments:
                documents[filename] = DocumentInfo(
                    filename=filename,
                    nb_pages=nb_pages,
                    nb_segments=len(segments),
                    nb_words=nb_words,
                    segment_start=offset,
                    segment_end=offset + len(segments),
                )
                all_segments.extend(segments)
                offset += len(segments)
                total_pages += nb_pages

        if not all_segments:
            return {"ok": False, "error": "No documents found or all failed to process."}

        # Embeddings (cache or compute)
        embeddings = None
        if not force_reindex and os.path.exists(cache_path):
            try:
                cached = np.load(cache_path, allow_pickle=False)
                if isinstance(cached, np.ndarray) and cached.shape[0] == len(all_segments):
                    embeddings = cached
            except Exception:
                logger.debug("Cache miss or corrupt: %s", cache_path)

        if embeddings is None:
            embeddings = self._encode(all_segments)
            try:
                os.makedirs(os.path.dirname(cache_path) or ".", exist_ok=True)
                np.save(cache_path, embeddings)
            except Exception:
                logger.warning("Failed to write embedding cache: %s", cache_path)

        # Build index
        with self._lock:
            self._segments = all_segments
            self._documents = documents
            self._build_index(embeddings)
            self._initialized = True

        elapsed = time.time() - start
        return {
            "ok": True,
            "files_indexed": len(documents),
            "total_segments": len(all_segments),
            "total_pages": total_pages,
            "total_words": sum(d.nb_words for d in documents.values()),
            "elapsed_seconds": round(elapsed, 2),
            "documents": {
                name: {
                    "pages": d.nb_pages,
                    "segments": d.nb_segments,
                    "words": d.nb_words,
                }
                for name, d in documents.items()
            },
        }

    def ingest_segments(
        self,
        segments: list[str],
        source_label: str = "chat",
    ) -> dict:
        """Append text segments to the live FAISS index without full re-index.

        Used by the chat extractor to incrementally add new messages.
        Returns a summary dict.
        """
        if not segments:
            return {"ok": True, "added_segments": 0, "total_segments": self.segment_count}

        start = time.time()

        # Encode new segments
        embeddings = self._encode(segments)

        with self._lock:
            # Initialise index if this is the first data
            if self._index is None:
                self._build_index(embeddings)
            else:
                self._index.add(embeddings)

            # Track as a virtual document
            offset = len(self._segments)
            self._segments.extend(segments)

            ts = time.strftime("%Y%m%d_%H%M%S")
            doc_name = f"_chat_{source_label}_{ts}"
            self._documents[doc_name] = DocumentInfo(
                filename=doc_name,
                nb_pages=1,
                nb_segments=len(segments),
                nb_words=sum(len(s.split()) for s in segments),
                segment_start=offset,
                segment_end=offset + len(segments),
            )
            self._initialized = True

        return {
            "ok": True,
            "added_segments": len(segments),
            "total_segments": len(self._segments),
            "source": source_label,
            "elapsed_seconds": round(time.time() - start, 2),
        }

    def search(
        self,
        query: str,
        top_k: int = 10,
        threshold: float = 0.8,
    ) -> list[SearchResult]:
        """Semantic search across indexed documents."""
        if not self.is_initialized:
            return []

        emb = self._encode([query])[0]
        distances, indices = self._index.search(np.array([emb]), k=min(top_k, len(self._segments)))

        results = []
        for i, idx in enumerate(indices[0]):
            if idx < 0 or idx >= len(self._segments):
                continue
            dist = float(distances[0][i])
            if dist > threshold:
                continue

            # Find which document this segment belongs to
            filename = "unknown"
            for name, doc in self._documents.items():
                if doc.segment_start <= idx < doc.segment_end:
                    filename = name
                    break

            results.append(SearchResult(
                filename=filename,
                segment_text=self._segments[idx],
                distance=round(dist, 4),
                segment_index=int(idx),
            ))

        return results

    def search_occurrences(self, keyword: str, top_k: int = 500) -> OccurrenceResult:
        """Count keyword occurrences using FAISS for candidate selection."""
        if not self.is_initialized:
            return OccurrenceResult(total=0, by_file={})

        emb = self._encode([keyword])[0]
        k = min(top_k, len(self._segments))
        distances, indices = self._index.search(np.array([emb]), k=k)

        by_file: dict[str, int] = {}
        total = 0

        for name, doc in self._documents.items():
            relevant = [
                self._segments[int(idx)]
                for i, idx in enumerate(indices[0])
                if doc.segment_start <= idx < doc.segment_end and distances[0][i] < 0.8
            ]
            if relevant:
                text = " ".join(relevant)
                count = len(re.findall(
                    r'\b' + re.escape(keyword) + r'\b', text, re.IGNORECASE
                ))
                if count > 0:
                    by_file[name] = count
                    total += count

        # Sort by count descending
        by_file = dict(sorted(by_file.items(), key=lambda x: x[1], reverse=True))
        return OccurrenceResult(total=total, by_file=by_file)

    def get_context(self, query: str, max_chars: int = 12000, top_k: int = 50) -> str:
        """Get relevant context for a query — designed to be passed to Copilot.

        Returns concatenated text segments relevant to the query,
        truncated to max_chars.
        """
        results = self.search(query, top_k=top_k, threshold=1.5)
        if not results:
            return ""

        context_parts: list[str] = []
        total_chars = 0
        seen_files: set[str] = set()

        for r in results:
            header = ""
            if r.filename not in seen_files:
                header = f"\n--- [{r.filename}] ---\n"
                seen_files.add(r.filename)

            chunk = header + r.segment_text
            if total_chars + len(chunk) > max_chars:
                remaining = max_chars - total_chars
                if remaining > 50:
                    context_parts.append(chunk[:remaining])
                break
            context_parts.append(chunk)
            total_chars += len(chunk)

        return "".join(context_parts)

    def list_documents(self, include_stats: bool = True) -> list[dict]:
        """List all indexed documents."""
        docs = []
        for name, info in self._documents.items():
            entry = {"filename": name}
            if include_stats:
                entry.update({
                    "pages": info.nb_pages,
                    "segments": info.nb_segments,
                    "words": info.nb_words,
                })
            docs.append(entry)
        return docs

    def session_brief(
        self,
        max_chars: int = 16_000,
        focus_query: Optional[str] = None,
    ) -> dict:
        """Build a comprehensive session briefing from the memory index.

        Runs multiple strategic queries to surface:
        - project overview (indexed documents)
        - key context from different angles
        - recent activity / chat history
        - pending tasks / TODOs

        Returns a dict with structured sections ready for the LLM.
        """
        start = time.time()

        # ── 1. Documents overview ──────────────────────────────────────
        docs = self.list_documents(include_stats=True)
        # Separate real docs from chat-injected virtual docs
        project_docs = [d for d in docs if not d["filename"].startswith("_chat_")]
        chat_docs = [d for d in docs if d["filename"].startswith("_chat_")]

        total_chat_segments = sum(d.get("segments", 0) for d in chat_docs)
        total_chat_words = sum(d.get("words", 0) for d in chat_docs)

        # ── 2. Multi-query retrieval (deduplicated) ───────────────────
        _QUERIES = [
            "project status progress current state",
            "TODO pending unfinished remaining tasks next steps",
            "recent activity latest changes updates new",
            "important decisions architecture choices design",
            "errors problems blockers issues bugs",
            "summary overview objectives goals",
        ]
        queries = list(_QUERIES)
        if focus_query:
            queries.insert(0, focus_query)

        seen_indices: set[int] = set()
        all_results: list[SearchResult] = []

        for q in queries:
            results = self.search(q, top_k=15, threshold=2.0)
            for r in results:
                if r.segment_index not in seen_indices:
                    seen_indices.add(r.segment_index)
                    all_results.append(r)

        # Sort by relevance (lowest distance first)
        all_results.sort(key=lambda r: r.distance)

        # ── 3. Build context string (budget-aware) ────────────────────
        context_parts: list[str] = []
        total_len = 0
        seen_files: set[str] = set()

        for r in all_results:
            header = ""
            if r.filename not in seen_files:
                header = f"\n--- [{r.filename}] ---\n"
                seen_files.add(r.filename)
            chunk = header + r.segment_text
            if total_len + len(chunk) > max_chars:
                remaining = max_chars - total_len
                if remaining > 50:
                    context_parts.append(chunk[:remaining])
                break
            context_parts.append(chunk)
            total_len += len(chunk)

        context_text = "".join(context_parts)

        # ── 4. Assemble brief ─────────────────────────────────────────
        return {
            "ok": True,
            "overview": {
                "total_documents": len(project_docs),
                "total_chat_sources": len(chat_docs),
                "total_segments": self.segment_count,
                "total_chat_segments": total_chat_segments,
                "total_chat_words": total_chat_words,
                "documents": [
                    {
                        "name": d["filename"],
                        "segments": d.get("segments", 0),
                        "words": d.get("words", 0),
                    }
                    for d in project_docs
                ],
            },
            "context": context_text,
            "context_chars": len(context_text),
            "queries_used": queries,
            "unique_segments_retrieved": len(all_results),
            "focus_query": focus_query,
            "elapsed_seconds": round(time.time() - start, 2),
        }

    def status(self) -> dict:
        """Return current engine status."""
        return {
            "initialized": self._initialized,
            "device": self.device,
            "model": self.model_name,
            "dimension": self.dimension,
            "document_count": self.document_count,
            "segment_count": self.segment_count,
            "segment_size": self.segment_size,
            "index_type": "FlatL2" if self._index else None,
        }

    def compact(
        self,
        max_segments: int = 500,
        keep_recent_hours: int = 24,
        strategy: str = "dedup_merge",
    ) -> dict:
        """Compact the in-memory index by removing near-duplicates and merging short segments.

        Strategies:
        - dedup_merge: remove near-duplicate segments (cosine distance < 0.05),
          then merge consecutive short segments from the same doc.
        - top_k: keep only the max_segments most central segments.

        Returns a summary with before/after counts.
        """
        if not self.is_initialized or not self._segments:
            return {"ok": True, "before": 0, "after": 0, "removed": 0, "strategy": strategy}

        start = time.time()
        before = len(self._segments)

        if before <= max_segments:
            return {
                "ok": True,
                "before": before,
                "after": before,
                "removed": 0,
                "strategy": strategy,
                "note": "Already under target — no compaction needed.",
            }

        faiss = _ensure_faiss()

        if strategy == "top_k":
            # Keep the max_segments segments closest to the centroid
            embeddings = self._encode(self._segments)
            centroid = embeddings.mean(axis=0, keepdims=True)
            distances = np.linalg.norm(embeddings - centroid, axis=1)
            top_indices = np.argsort(distances)[:max_segments]
            top_indices = sorted(top_indices)  # preserve order

            new_segments = [self._segments[i] for i in top_indices]
            new_embeddings = embeddings[top_indices]

        else:  # dedup_merge
            # 1. Remove near-duplicates via cosine similarity matrix (vectorized)
            embeddings = self._encode(self._segments)
            n = len(self._segments)
            # Compute pairwise L2 distances efficiently via broadcasting
            # For normalized embeddings: dist_sq = 2 - 2*dot(a,b)
            # Use chunked dot product to limit memory for very large n
            keep = [True] * n
            if n <= 5000:
                # Vectorized approach: compute full similarity matrix at once
                sim = embeddings @ embeddings.T  # n×n cosine similarity
                for i in range(n):
                    if not keep[i]:
                        continue
                    for j in range(i + 1, n):
                        if not keep[j]:
                            continue
                        dist = 2.0 - 2.0 * float(sim[i, j])
                        if dist < 0.05:
                            if len(self._segments[j]) > len(self._segments[i]):
                                keep[i] = False
                                break
                            else:
                                keep[j] = False
            else:
                # For very large indices, process in chunks
                chunk_size = 1000
                for start in range(0, n, chunk_size):
                    end = min(start + chunk_size, n)
                    chunk_emb = embeddings[start:end]
                    sim_chunk = chunk_emb @ embeddings.T
                    for local_i in range(end - start):
                        i = start + local_i
                        if not keep[i]:
                            continue
                        for j in range(i + 1, n):
                            if not keep[j]:
                                continue
                            dist = 2.0 - 2.0 * float(sim_chunk[local_i, j])
                            if dist < 0.05:
                                if len(self._segments[j]) > len(self._segments[i]):
                                    keep[i] = False
                                    break
                                else:
                                    keep[j] = False

            deduped = [(self._segments[i], embeddings[i]) for i in range(len(self._segments)) if keep[i]]

            # 2. Merge consecutive short segments (< 64 chars)
            merged_segments: list[str] = []
            merged_embeddings: list[np.ndarray] = []
            buf = ""
            for seg, emb in deduped:
                if len(seg) < 64 and buf:
                    buf += " " + seg
                else:
                    if buf:
                        merged_segments.append(buf)
                        merged_embeddings.append(self._encode([buf])[0])
                    buf = seg
            if buf:
                merged_segments.append(buf)
                merged_embeddings.append(self._encode([buf])[0])

            # 3. If still over target, keep top-k by centroid distance
            if len(merged_segments) > max_segments:
                embs = np.array(merged_embeddings)
                centroid = embs.mean(axis=0, keepdims=True)
                distances = np.linalg.norm(embs - centroid, axis=1)
                top_indices = np.argsort(distances)[:max_segments]
                top_indices = sorted(top_indices)
                new_segments = [merged_segments[i] for i in top_indices]
                new_embeddings = np.array([merged_embeddings[i] for i in top_indices])
            else:
                new_segments = merged_segments
                new_embeddings = np.array(merged_embeddings)

        # Rebuild index
        with self._lock:
            self._segments = new_segments
            self._documents = {
                "_compacted": DocumentInfo(
                    filename="_compacted",
                    nb_pages=1,
                    nb_segments=len(new_segments),
                    nb_words=sum(len(s.split()) for s in new_segments),
                    segment_start=0,
                    segment_end=len(new_segments),
                )
            }
            self._build_index(new_embeddings.astype("float32"))

        return {
            "ok": True,
            "before": before,
            "after": len(new_segments),
            "removed": before - len(new_segments),
            "strategy": strategy,
            "elapsed_seconds": round(time.time() - start, 2),
        }

    def get_segment_text(self, doc_name: str) -> str | None:
        """Return concatenated text for a document (for MCP resources)."""
        doc = self._documents.get(doc_name)
        if doc is None:
            return None
        return "\n".join(self._segments[doc.segment_start:doc.segment_end])

    def transcribe(self, file_path: str, language: str = "fr") -> dict:
        """Transcribe an audio file using Whisper."""
        if not os.path.isfile(file_path):
            return {"ok": False, "error": f"File not found: {file_path}"}

        ext = os.path.splitext(file_path)[1].lower()
        if ext not in {".mp3", ".wav", ".ogg", ".flac"}:
            return {"ok": False, "error": f"Unsupported audio format: {ext}"}

        try:
            start = time.time()
            text, _ = _extract_audio(file_path, language=language)
            return {
                "ok": True,
                "text": text,
                "file": os.path.basename(file_path),
                "language": language,
                "elapsed_seconds": round(time.time() - start, 2),
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}
